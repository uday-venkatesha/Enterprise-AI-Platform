import io

from pptx import Presentation


def extract_text_from_pptx(raw_bytes: bytes) -> str:
    presentation = Presentation(io.BytesIO(raw_bytes))

    lines: list[str] = []
    # A presentation is a list of slides...
    for slide in presentation.slides:
        # ...and each slide holds SHAPES (text boxes, titles, etc.). The text
        # isn't on the slide directly — it's inside shapes. Not every shape has
        # text (an image doesn't), so we check has_text_frame first.
        for shape in slide.shapes:
            if shape.has_text_frame:
                # A text frame is itself made of paragraphs — pull each one.
                for paragraph in shape.text_frame.paragraphs:
                    text = "".join(run.text for run in paragraph.runs)
                    if text.strip():
                        lines.append(text)

    return "\n".join(lines).strip()